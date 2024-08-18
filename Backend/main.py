from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from transformer_model import get_summarizer, summarize_text
import io
import PyPDF2
import docx
import pptx

app = FastAPI()

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Your React app's URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize the summarizer
summarizer = get_summarizer()

class SummarizeRequest(BaseModel):
    text: str
    mode: str = 'Paragraph'
    length: str = 'Concise'

@app.post("/summarize")
async def summarize(request: SummarizeRequest):
    try:
        # Set max_length and min_length based on length
        if request.length == 'Concise':
            max_length, min_length = 100, 30
        else:  # Standard
            max_length, min_length = 200, 50

        summary = summarize_text(summarizer, request.text, max_length, min_length, request.mode)
        return {"summary": summary}
    except Exception as e:
        print(f"Error in summarization: {str(e)}")
        return JSONResponse(status_code=500, content={"detail": f"Error in summarization: {str(e)}"})

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    content = await file.read()
    file_extension = file.filename.split('.')[-1].lower()

    if file_extension == 'txt':
        text = content.decode('utf-8')
    elif file_extension == 'pdf':
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
    elif file_extension == 'docx':
        doc = docx.Document(io.BytesIO(content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    elif file_extension == 'pptx':
        prs = pptx.Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
    else:
        raise HTTPException(status_code=500, detail="Unsupported file format")

    text_split_line = text.split("\n")
    word_count = 0
    
    for line in text_split_line:
        word_count += len(line.strip().split(" "))
    
    if word_count > 2000:
        raise HTTPException(status_code=500, detail="Word limit exceeded")

    return {"text": text, 
            "count": word_count
        }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)